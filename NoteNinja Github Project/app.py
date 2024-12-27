from flask import Flask, render_template,request, g, send_from_directory, session  #type: ignore
from openai import OpenAI #type: ignore
from dotenv import load_dotenv
import os
import re
import io  
from docx import Document
from pptx import Presentation
import PyPDF2
from PyPDF2 import PdfReader
import json
import jsonify
from supabase import create_client



app = Flask(__name__)
# app.secret_key = os.getenv("session_key")

load_dotenv()

# open ai setup
client = OpenAI(
        api_key = os.getenv('api_key'),
)

# supabase setup
url = os.getenv("SUPABASE_URL")
key = os.getenv("SUPABASE_KEY")

supabase = create_client(url, key)



# this will make sure that application/js is enforced when i try to import functions / modules
@app.route('/static/<path:filename>')
def static_files(filename):
    return send_from_directory('static', filename, mimetype='application/javascript')

# main page where the user will upload their files
@app.route("/", methods=["POST", "GET"])
def home():
        

    return render_template("index.html")

# simple about us page
@app.route("/about")
def about():
    return render_template("about.html");

# simple contact us page 
@app.route("/contact")
def contact():
    return render_template("contact.html")

# this is where everything will be processed and renderd
@app.route("/notes", methods=["POST", "GET"])
def notes():
    if request.method == "POST":
        # where the chatgpt api will work
        storeFile = request.files["file-upload"]
        file_name = storeFile.filename
        file_extension = file_name.lower().split('.')[-1]
        
        # Read file content based on file type
        file_contents = ""
        file_bytes = io.BytesIO(storeFile.read())
        
        # reads the file based on the file type it is..
        try:
            if file_extension == 'txt':
                file_bytes.seek(0)
                file_contents = file_bytes.read().decode('utf-8') #converts file into readbale utf-8
                
            elif file_extension == 'pdf':
                pdf_reader = PyPDF2.PdfReader(file_bytes)
                for page in pdf_reader.pages:
                    file_contents += page.extract_text()
                    
            elif file_extension == 'docx':
                doc = Document(file_bytes)
                for para in doc.paragraphs:
                    file_contents += para.text + '\n'
                    
            elif file_extension == 'pptx':
                prs = Presentation(file_bytes)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            file_contents += shape.text + '\n'
            else:
                return render_template("index.html", error="Unsupported file type")
            
            # Add character limit check (approximately 4000 tokens ≈ 16000 characters)
            MAX_CHARS = 16000
            if len(file_contents) > MAX_CHARS:
                truncated_message = (
                    f"⚠️ Note: The document was too long and has been truncated to the first {MAX_CHARS} characters for processing.\n\n"
                )
                file_contents = file_contents[:MAX_CHARS] + "\n...[truncated]"
            else:
                truncated_message = ""
                
            completion = client.chat.completions.create(
                model="gpt-4",
                messages = [
                # clear instructions to the model that tells it what to do. I used compleation because i realized that assistants was for chatbots and not for this
                {"role": "system", "content": "You are a diligent and insightful assistant dedicated to transforming documents into well-structured, concise, and informative notes. The notes should be split into different sections with a clear title "},
                 # user input which will be the file uploaded through html form
                {'role': 'user', 'content': file_contents}
                    
                ],
                
                temperature=1,
                max_tokens=1000
            )
            
            response = completion.choices[0].message.content
            # Add truncation message to response if needed
            if truncated_message:
                response = truncated_message + response
        
            # splits the response into a list of notes
            notes = response.split("\n\n")
            print(notes)
            
            # save notes into a database
     

            
            
            return render_template("notes.html", notes=notes)
        # just incase the file is not readable    
        except Exception as e:
            return render_template("index.html", error=f"Error processing file: {str(e)}")
    return render_template("notes.html")

# where said user can view the notes they have uploaded



@app.route('/oldnotes')
def oldnotes():
    return render_template("oldnotes.html")

@app.route("/savednotes")
def savedNotesInner():
    return render_template("savednotes.html")

    

app.run(host="0.0.0.0", port=5000, debug=True)